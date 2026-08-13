#!/usr/bin/env python3
"""
Generate VAL CoPilot architecture + process-flow PowerPoint for the
Synthetic Contract Intelligence Tool-Layer POC.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

DOCS = Path(__file__).resolve().parent
PPTX_PATH = DOCS / "VAL_CoPilot_Architecture_and_Process_Flow.pptx"
ARTIFACT = Path("/opt/cursor/artifacts/VAL_CoPilot_Architecture_and_Process_Flow.pptx")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x0B, 0x1F, 0x33)
TEAL = RGBColor(0x0E, 0x6B, 0x6B)
TEAL_LIGHT = RGBColor(0xD9, 0xF3, 0xF3)
SLATE = RGBColor(0x33, 0x41, 0x55)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xC2, 0x41, 0x0C)
ORANGE_LIGHT = RGBColor(0xFF, 0xED, 0xD5)
BLUE = RGBColor(0x1D, 0x4E, 0xD8)
BLUE_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)
PURPLE = RGBColor(0x6D, 0x28, 0xD9)
PURPLE_LIGHT = RGBColor(0xED, 0xE9, 0xFE)
GREEN = RGBColor(0x04, 0x78, 0x57)
GREEN_LIGHT = RGBColor(0xD1, 0xFA, 0xE5)
GOLD = RGBColor(0xB4, 0x53, 0x09)
GOLD_LIGHT = RGBColor(0xFE, 0xF3, 0xC7)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
RED = RGBColor(0xB9, 0x1C, 0x1C)
RED_LIGHT = RGBColor(0xFE, 0xE2, 0xE2)


def _set_run(run, text: str, size: int, bold: bool = False, color: RGBColor = SLATE) -> None:
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _fill_shape(shape, fill: RGBColor, line: RGBColor, line_pt: float = 1.25) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_pt)


def add_box(slide, left, top, width, height, fill, line, title, subtitle="", title_size=12, sub_size=10, title_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _fill_shape(shape, fill, line)
    try:
        shape.adjustments[0] = 0.12
    except Exception:  # noqa: BLE001
        pass
    tf = shape.text_frame
    tf.word_wrap = True
    shape.text_frame.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    _set_run(run, title, title_size, bold=True, color=title_color or NAVY)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        _set_run(run2, subtitle, sub_size, color=SLATE)
    return shape


def add_title_bar(slide, text: str, subtitle: str = "") -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.85))
    _fill_shape(bar, NAVY, NAVY, 0)
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.55))
    run = tb.text_frame.paragraphs[0].add_run()
    _set_run(run, text, 20, bold=True, color=WHITE)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.4), Inches(0.9), Inches(12.5), Inches(0.35))
        run2 = sub.text_frame.paragraphs[0].add_run()
        _set_run(run2, subtitle, 12, color=SLATE)


def add_arrow_right(slide, left, top, width=Inches(0.3), height=Inches(0.18), color=GRAY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    _fill_shape(shape, color, color, 0)
    return shape


def add_footer(slide, page: str) -> None:
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    _set_run(run, f"Synthetic Contract Intelligence Tool-Layer POC · {page}", 9, color=GRAY)


def slide_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _fill_shape(bg, LIGHT, LIGHT, 0)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.0), SLIDE_W, Inches(2.8))
    _fill_shape(accent, NAVY, NAVY, 0)
    title = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.7), Inches(0.8))
    p = title.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), "VAL CoPilot", 34, bold=True, color=WHITE)
    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.15), Inches(11.7), Inches(1.0))
    p2 = sub.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    _set_run(p2.add_run(), "Synthetic Contract Intelligence\nTool-Layer POC — Architecture & Process Flow", 22, bold=True, color=TEAL_LIGHT)
    note = slide.shapes.add_textbox(Inches(1.2), Inches(5.3), Inches(10.9), Inches(1.4))
    for i, line in enumerate([
        "Current: Streamlit → Flask /api/messages → LangChain or offline router → FastMCP → LinkSquares fixtures + persona memory",
        "Future: Streamlit/Teams → Azure AI Foundry Agent → same MCP/API tools → live LinkSquares later",
        "Compare hard-stop when IDs missing · Invoice/actual-spend OOS · Not Foundry-first today",
    ]):
        para = note.text_frame.paragraphs[0] if i == 0 else note.text_frame.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        _set_run(para.add_run(), line, 13, color=SLATE)


def slide_current_vs_future(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "1. Current vs Future Architecture", "Keep the tool layer stable; swap the cognitive host later")
    add_box(
        slide, Inches(0.4), Inches(1.4), Inches(6.1), Inches(5.2), TEAL_LIGHT, TEAL,
        "CURRENT — Tool-Layer POC",
        "Streamlit UI + Saved searches\n→ Flask POST /api/messages\n→ LangChain AgentExecutor\n   or offline cognitive router\n→ FastMCP tool server\n→ ContractRepository\n→ LinkSquares offline fixtures\n→ SQLite persona memory",
        16, 13,
    )
    add_box(
        slide, Inches(6.8), Inches(1.4), Inches(6.1), Inches(5.2), PURPLE_LIGHT, PURPLE,
        "FUTURE — Foundry-hosted agent",
        "Streamlit or Teams UI\n→ Azure AI Foundry Agent\n→ Same MCP / API tools\n→ ContractRepository\n→ Synthetic contract data now\n→ LinkSquares / CLM later\n→ Grounded answers",
        16, 13,
    )
    add_footer(slide, "Slide 2")


def slide_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(
        slide,
        "2. Tool-Layer System Architecture",
        "Four layers + persona memory + replaceable ContractRepository",
    )
    add_box(slide, Inches(0.35), Inches(1.2), Inches(12.6), Inches(1.05), BLUE_LIGHT, BLUE,
            "Layer 1 — Validation UI (test_ui)",
            "Streamlit chat · persona picker · Saved searches sidebar → COPILOT_MESSAGES_URL", 13, 11)
    add_box(slide, Inches(0.35), Inches(2.35), Inches(12.6), Inches(1.25), TEAL_LIGHT, TEAL,
            "Layer 2 — Cognitive Routing + Memory (copilot_agent + memory/)",
            "Flask /api/messages · /api/memory/* · LangChain / offline_router · SYSTEM_PROMPT · "
            "invoice OOS · compare hard-stop · SQLite persona store",
            13, 11)
    add_box(slide, Inches(0.35), Inches(3.75), Inches(12.6), Inches(1.55), GREEN_LIGHT, GREEN,
            "Layer 3 — Data Retrieval (mcp_server)",
            "ContractRepository → analytics / risk helpers\n"
            "Tools: search · profile · compare (N-way) · missing-fields · expiring · spend rollup · "
            "overlaps · risk · blob search · health",
            13, 11)
    add_box(slide, Inches(0.35), Inches(5.45), Inches(12.6), Inches(1.15), GOLD_LIGHT, GOLD,
            "Layer 4 — Synthetic Gold (LinkSquares fixtures)",
            "LinSquare_Contracts_100_Updated_30bb.json + agreement_9a06.json · future live Fabric / LinkSquares",
            13, 11)
    add_footer(slide, "Slide 3")


def slide_request_flow(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "3. End-to-End Request Process Flow", "Deterministic offline path used for demo reliability")
    steps = [
        ("1. User prompt", "Streamlit chat", BLUE_LIGHT, BLUE, 0.35),
        ("2. Guardrails", "OOS / hard-stop / memory", RED_LIGHT, RED, 2.9),
        ("3. Route intent", "offline_router / LLM", TEAL_LIGHT, TEAL, 5.45),
        ("4. MCP tools", "ContractRepository", GREEN_LIGHT, GREEN, 8.0),
        ("5. Reply + save", "Markdown · auto-save", ORANGE_LIGHT, ORANGE, 10.55),
    ]
    for title, sub, fill, line, x in steps:
        add_box(slide, Inches(x), Inches(1.55), Inches(2.35), Inches(1.25), fill, line, title, sub, 12, 10)
    for x in (2.55, 5.1, 7.65, 10.2):
        add_arrow_right(slide, Inches(x), Inches(2.05), color=GRAY)

    add_box(
        slide, Inches(0.4), Inches(3.05), Inches(12.5), Inches(3.55), LIGHT, GRAY,
        "Intent → primary tool / outcome",
        "Show contracts for vendor → search_contracts\n"
        "Details for CON-0002 → get_contract_profile\n"
        "Overlapping contracts → find_overlaps\n"
        "Unusual payment / high rates → explain_contract_risk\n"
        "Compare CON-0001 vs CON-0002 → compare_contracts (N-way any IDs)\n"
        "Compare unknown vendor/ID → hard-stop message only (no table / no recommendation)\n"
        "Need action next 90 days → get_expiring_contracts\n"
        "Previous / saved searches → persona memory recall (/api/memory/*)\n"
        "Invoice spend also? → exact OOS refusal (no tools)",
        13, 11,
    )
    add_footer(slide, "Slide 4")


def slide_hardstop_memory(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(
        slide,
        "3b. Compare Hard-Stop & Persona Memory",
        "Deterministic offline paths — no LLM default pairs, no invented history",
    )
    add_box(
        slide, Inches(0.4), Inches(1.35), Inches(6.2), Inches(5.3), RED_LIGHT, RED,
        "Compare hard-stop",
        "If any requested supplier or contract ID\n"
        "cannot be resolved:\n\n"
        "Return ONLY:\n"
        "“The contract information requested\n"
        "for the comparison is not available\n"
        "at the moment”\n\n"
        "• No comparative table\n"
        "• No ## Recommendation\n"
        "• No CON-0001 vs CON-0002 default\n"
        "• Compare intents → offline_router",
        15, 12,
    )
    add_box(
        slide, Inches(6.9), Inches(1.35), Inches(6.0), Inches(5.3), TEAL_LIGHT, TEAL,
        "Persona memory",
        "SQLite: data/persona_memory.sqlite\n"
        "(override VAL_MEMORY_DB)\n\n"
        "• Auto-save search-like queries\n"
        "• Pin / filter / delete in sidebar\n"
        "• GET/POST /api/memory/searches\n"
        "• GET /api/memory/recall\n"
        "• Memory-recall → offline store\n"
        "• Streamlit owns writes when\n"
        "  clientPersistsMemory is set\n\n"
        "memory/store.py + test_ui sidebar",
        15, 12,
    )
    add_footer(slide, "Slide 5")


def slide_repository(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "4. ContractRepository Abstraction", "Data source is replaceable without changing tool contracts")
    add_box(slide, Inches(0.5), Inches(1.4), Inches(4.0), Inches(5.2), TEAL_LIGHT, TEAL,
            "ContractRepository",
            "list_all()\nget_by_id(contract_id)\nsearch(filters)\nget_by_vendor(vendor)",
            16, 14)
    add_arrow_right(slide, Inches(4.7), Inches(3.7), Inches(0.5), Inches(0.28), GRAY)
    add_box(slide, Inches(5.4), Inches(1.4), Inches(3.6), Inches(2.3), GREEN_LIGHT, GREEN,
            "Today", "FabricContractRepository\n→ LinkSquares fixtures\n→ Gold-shaped projections", 15, 13)
    add_box(slide, Inches(5.4), Inches(4.0), Inches(3.6), Inches(2.6), PURPLE_LIGHT, PURPLE,
            "Later", "Live LinkSquares / CLM\nsame tool signatures\nsame snake_case projections", 15, 13)
    add_box(slide, Inches(9.3), Inches(1.4), Inches(3.5), Inches(5.2), GOLD_LIGHT, GOLD,
            "Consumers", "search_contracts\nget_contract_profile\nfind_overlaps\nexplain_contract_risk\ncompare / missing / expiring", 15, 12)
    add_footer(slide, "Slide 6")


def slide_new_tools(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "5. Structured Intelligence Tools (P0/P1)", "Grounded metadata tools — not document search")
    tools = [
        ("search_contracts", "Vendor/BU/status/type filters\nStable snake_case catalog rows", GREEN_LIGHT, GREEN, 0.4, 1.3),
        ("get_contract_profile", "One contract_id profile\n+ missing_fields", TEAL_LIGHT, TEAL, 4.6, 1.3),
        ("compare_contracts", "Pairwise / N-way any IDs\nHard-stop if unresolved", ORANGE_LIGHT, ORANGE, 8.8, 1.3),
        ("explain_contract_risk", "known_facts vs computed_risks\nmissing_data + review action", PURPLE_LIGHT, PURPLE, 0.4, 4.0),
        ("find_overlaps", "Same-vendor date overlaps\nOverlapFlag-aware pairs", BLUE_LIGHT, BLUE, 4.6, 4.0),
        ("Demo fixtures", "LinkSquares 100 contracts\nMicrosoft/AWS overlaps & Net 180", GOLD_LIGHT, GOLD, 8.8, 4.0),
    ]
    for title, sub, fill, line, x, y in tools:
        add_box(slide, Inches(x), Inches(y), Inches(3.9), Inches(2.3), fill, line, title, sub, 14, 12)
    add_footer(slide, "Slide 7")


def slide_risk_flow(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "6. Overlap & Risk Explanation Process Flow", "Separate facts from computed risks — no invented findings")
    add_box(slide, Inches(0.4), Inches(1.4), Inches(4.0), Inches(5.2), ORANGE_LIGHT, ORANGE,
            "find_overlaps",
            "1. Load contracts via repository\n2. Group by vendor\n3. Prefer OverlapFlag=Yes pairs\n4. Emit overlap_start/end\n5. why_flagged + source\n\nDemo: CON-0024 ∩ CON-0029",
            15, 12)
    add_box(slide, Inches(4.7), Inches(1.4), Inches(4.0), Inches(5.2), PURPLE_LIGHT, PURPLE,
            "explain_contract_risk",
            "Flags only if data supports:\n• missing renewal\n• missing rate card\n• expiring soon\n• Net 90+ terms\n• high ACV outlier\n• high supplier risk\n• overlapping contract",
            15, 12)
    add_box(slide, Inches(9.0), Inches(1.4), Inches(3.9), Inches(5.2), TEAL_LIGHT, TEAL,
            "Response sections",
            "known_facts\ncomputed_risks\nmissing_data\nrecommended_review_action\nsource\n\nNever invent unsupported risks",
            15, 12)
    add_footer(slide, "Slide 8")


def slide_tool_inventory(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "7. Full MCP Tool Inventory", "fabric_data FastMCP server")
    left = [
        "search_contracts",
        "get_contract_profile",
        "find_overlaps",
        "explain_contract_risk",
        "compare_contracts (N-way; hard-stop)",
    ]
    right = [
        "check_missing_contract_fields",
        "get_expiring_contracts",
        "get_vendor_spend_summary (ACV rollup only)",
        "search_cloud_blob_contracts (docs)",
        "fabric_health_check",
    ]
    add_box(slide, Inches(0.5), Inches(1.4), Inches(6.0), Inches(4.8), GREEN_LIGHT, GREEN,
            "Structured / analytics", "\n".join(left), 15, 13)
    add_box(slide, Inches(6.9), Inches(1.4), Inches(5.9), Inches(4.8), BLUE_LIGHT, BLUE,
            "Commercial / docs / health", "\n".join(right), 15, 13)
    add_footer(slide, "Slide 9")


def slide_foundry_future(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(
        slide,
        "8. Optional Future — Azure AI Foundry Deploy",
        "deploy_to_foundry.py is optional; current POC does not require Foundry",
    )
    steps = [
        ("Root .env", "AZURE_FOUNDRY_\nCONNECTION_STRING", TEAL_LIGHT, TEAL, 0.4),
        ("Entra auth", "DefaultAzureCredential", ORANGE_LIGHT, ORANGE, 3.0),
        ("Wrap tools", "FunctionTool ToolSet\n(7 MCP tools)", GREEN_LIGHT, GREEN, 5.6),
        ("create_agent", "SYSTEM_PROMPT\nhard-stop + OOS", BLUE_LIGHT, BLUE, 8.2),
        ("Future host", "Teams / Foundry UI\nextend ToolSet", PURPLE_LIGHT, PURPLE, 10.8),
    ]
    for title, sub, fill, line, x in steps:
        add_box(slide, Inches(x), Inches(2.0), Inches(2.3), Inches(2.2), fill, line, title, sub.replace("\n", " "), 13, 11)
    for x in (2.75, 5.35, 7.95, 10.55):
        add_arrow_right(slide, Inches(x), Inches(3.0), Inches(0.22), Inches(0.2), GRAY)
    add_box(
        slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(1.8), LIGHT, GRAY,
        "Deploy script ToolSet today (7 tools) — extend next to full MCP inventory",
        "search_contracts · get_contract_profile · get_expiring_contracts · get_vendor_spend_summary · "
        "find_overlaps · explain_contract_risk · search_cloud_blob_contracts\n"
        "Local FastMCP also has compare_contracts + check_missing_contract_fields. "
        "SYSTEM_PROMPT encodes compare hard-stop + persona-memory guidance.",
        13, 11,
    )
    add_footer(slide, "Slide 10")


def slide_summary(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "9. Demo-Ready Slice Summary", "Reliable · safe · synthetic")
    items = [
        ("Safe", "Invoice OOS + compare hard-stop (exact messages; no default CON-* pair)", RED, RED_LIGHT),
        ("Structured tools", "search / profile / N-way compare / overlaps / risk via ContractRepository", GREEN, GREEN_LIGHT),
        ("Persona memory", "SQLite save / pin / filter / recall previous searches per persona", TEAL, TEAL_LIGHT),
        ("LinkSquares fixtures", "CON-* offline data + docs/demo_script.md; live CLM later", PURPLE, PURPLE_LIGHT),
        ("Foundry optional", "Flask + offline router today; Foundry ToolSet wraps 7 tools", GOLD, GOLD_LIGHT),
    ]
    for i, (title, body, color, light) in enumerate(items):
        y = 1.25 + i * 1.05
        add_box(slide, Inches(0.5), Inches(y), Inches(2.6), Inches(0.85), light, color, title, title_size=14, title_color=color)
        add_box(slide, Inches(3.3), Inches(y), Inches(9.5), Inches(0.85), WHITE, color, body, title_size=13)
    add_footer(slide, "Slide 11")


def build_pptx(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_cover(prs)
    slide_current_vs_future(prs)
    slide_architecture(prs)
    slide_request_flow(prs)
    slide_hardstop_memory(prs)
    slide_repository(prs)
    slide_new_tools(prs)
    slide_risk_flow(prs)
    slide_tool_inventory(prs)
    slide_foundry_future(prs)
    slide_summary(prs)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))


def main() -> None:
    build_pptx(PPTX_PATH)
    ARTIFACT.parent.mkdir(parents=True, exist_ok=True)
    ARTIFACT.write_bytes(PPTX_PATH.read_bytes())
    print(f"Wrote {PPTX_PATH} ({PPTX_PATH.stat().st_size} bytes)")
    print(f"Wrote {ARTIFACT} ({ARTIFACT.stat().st_size} bytes)")


if __name__ == "__main__":
    main()
